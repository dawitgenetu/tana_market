from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Define Brand Colors
    PRIMARY_COLOR = RGBColor(79, 70, 229)  # Indigo/Purple-ish (modern web app feel)
    SECONDARY_COLOR = RGBColor(30, 41, 59) # Dark Slate
    ACCENT_COLOR = RGBColor(245, 158, 11)  # Amber for highlights
    TEXT_COLOR = RGBColor(51, 65, 85)      # Slate for text
    WHITE = RGBColor(255, 255, 255)
    LIGHT_BG = RGBColor(248, 250, 252)     # Very light gray

    def apply_slide_design(slide, title_text):
        # 1. Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = LIGHT_BG

        # 2. Header Strip
        left = top = Inches(0)
        width = prs.slide_width
        height = Inches(1.2)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY_COLOR
        shape.line.fill.background() # No border

        # 3. Add Logo/Icon placeholder (Circle)
        logo_dia = Inches(0.8)
        logo_left = Inches(0.5)
        logo_top = Inches(0.2)
        logo = slide.shapes.add_shape(MSO_SHAPE.OVAL, logo_left, logo_top, logo_dia, logo_dia)
        logo.fill.solid()
        logo.fill.fore_color.rgb = WHITE
        logo.line.color.rgb = ACCENT_COLOR
        
        # 4. Title Text
        text_left = Inches(1.5)
        text_top = Inches(0.3)
        text_width = Inches(8)
        text_height = Inches(0.8)
        
        textbox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.LEFT

        # 5. Footer Line
        footer_height = Inches(0.05)
        footer_top = prs.slide_height - Inches(0.5)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), footer_top, prs.slide_width, footer_height)
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_COLOR
        line.line.fill.background()

    def add_content_slide(title, content_points, layout_idx=6): # 6 is usually Blank, we build custom
        slide = prs.slides.add_slide(prs.slide_layouts[6]) 
        apply_slide_design(slide, title)
        
        # Content box
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8.5)
        height = Inches(5.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for point in content_points:
            p = tf.add_paragraph()
            p.text = "\u2022 " + point # Bullet point
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_COLOR
            p.font.name = 'Calibri'
            p.space_after = Pt(14)
            p.space_before = Pt(0)
            p.level = 0

    def add_title_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background Image/Color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY_COLOR
        
        # Title
        title_top = Inches(2.5)
        title = slide.shapes.add_textbox(Inches(1), title_top, Inches(8.5), Inches(2))
        title_tf = title.text_frame
        p = title_tf.paragraphs[0]
        p.text = "TANA MARKET"
        p.font.size = Pt(64)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        sub_top = Inches(4)
        sub = slide.shapes.add_textbox(Inches(1), sub_top, Inches(8.5), Inches(1))
        sub_tf = sub.text_frame
        p = sub_tf.add_paragraph()
        p.text = "Comprehensive E-Commerce Platform Documentation"
        p.font.size = Pt(28)
        p.font.color.rgb = ACCENT_COLOR
        p.alignment = PP_ALIGN.CENTER

        # Presenter Name
        name_top = Inches(5.5)
        name = slide.shapes.add_textbox(Inches(1), name_top, Inches(8.5), Inches(1))
        name_tf = name.text_frame
        p = name_tf.add_paragraph()
        p.text = "Presented by: Tana Market Team\nDate: January 2026"
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # --- SLIDES GENERATION (18 Slides) ---

    # 1. Title Slide
    add_title_slide()

    # 2. Introduction
    add_content_slide("Introduction", [
        "Tana Market is a modern, full-stack e-commerce platform facilitating seamless interaction between vendors and customers.",
        "Built to address the digital gap in local retail markets by providing a robust online presence.",
        "Integrates advanced features like Role-Based Access Control (RBAC), Real-time Tracking, and Secure Digital Payments."
    ])

    # 3. Statement of the Problem
    add_content_slide("Statement of the Problem", [
        "Limited Market Reach: Local vendors struggle to reach a wider audience beyond their physical location.",
        "Inefficient Inventory Management: Manual tracking leads to stock discrepancies and operational losses.",
        "Lack of Trust: Customers often face issues with payment security and order transparency.",
        "Fragmented Processes: Disconnected systems for sales, shipping, and payments cause delays."
    ])

    # 4. General Objective
    add_content_slide("General Objective", [
        "To develop a scalable, secure, and user-friendly web-based e-commerce application that streamlines the entire online shopping lifecycle.",
        "To bridge the gap between traditional retail and digital commerce in the local market.",
        "To provide a centralized platform for managing products, orders, users, and financial transactions efficiently."
    ])

    # 5. Specific Objectives (1/2)
    add_content_slide("Specific Objectives (1/2)", [
        "Develop a responsive user interface for Customers, Managers, and Admins using React.",
        "Implement a Role-Based Access Control (RBAC) system to ensure data security and appropriate user permissions.",
        "Create a robust product management system allowing dynamic categorization, image uploads, and inventory tracking.",
        "Integrate Chapa payment gateway for secure, real-time local currency transactions."
    ])

    # 6. Specific Objectives (2/2)
    add_content_slide("Specific Objectives (2/2)", [
        "Establish an automated order tracking system with unique tracking IDs (e.g., TANA-2026-XXXX).",
        "Implement real-time dashboards for Admins to monitor sales statistics and user activities.",
        "Provide a feedback loop through a product rating and review system to enhance trust.",
        "Ensure data integrity and security through encryption and secure API authentication."
    ])

    # 7. Scope of the Project
    add_content_slide("Project Scope (In-Scope)", [
        "User Registration & Authentication (Login/Register/Logout).",
        "Product Catalog: Search, Filter, Sort, and View Details.",
        "Shopping Cart & Checkout Process.",
        "Order Management: lifecycle from 'Pending' to 'Delivered'.",
        "Admin Dashboard: User audits, Reports, and System Settings.",
        "Manager Dashboard: Product CRUD and Order Processing."
    ])

    # 8. Scope & Limitations
    add_content_slide("Scope Limitations", [
        "Mobile App: The current phase focuses only on the Web Application (Responsive), not a native mobile app.",
        "International Shipping: Initial launch supports local logistics only (e.g., Bahir Dar).",
        "Offline Mode: The system requires an active internet connection to function.",
        "Multi-Vendor Marketplace: Currently operates as a single-store inventory model, not a multi-vendor platform."
    ])

    # 9. System Overview
    add_content_slide("System Architecture", [
        "Frontend: Single Page Application (SPA) built with React + Vite.",
        "Backend: RESTful API built with Node.js and Express.",
        "Database: MongoDB for flexible, document-oriented data storage.",
        "Deployment: Docker-ready, cloud-agnostic architecture."
    ])

    # 10. Technology Stack
    add_content_slide("Technology Stack", [
        "Frontend: React.js, TailwindCSS (Styling), Framer Motion (Animations), Axios (API).",
        "Backend: Node.js, Express.js, JWT (Auth), Multer (File Uploads).",
        "Database: MongoDB, Mongoose ODM.",
        "Tools: Git, Postman, VS Code, PlantUML (Documentation)."
    ])

    # 11. User Management (Module)
    add_content_slide("Module: User Management", [
        "Registration: Secure signup with email validation.",
        "Authentication: JWT (JSON Web Tokens) for session management.",
        "RBAC Implementation:",
        "   - Admin: Full system control.",
        "   - Manager: Inventory & Order focus.",
        "   - Customer: Shopping & Profile focus."
    ])

    # 12. Product Management (Module)
    add_content_slide("Module: Product & Inventory", [
        "CRUD Operations: Create, Read, Update, Delete products.",
        "Categorization: Hierarchical organization (Electronics, Fashion, etc.).",
        "Stock Control: Auto-decrement stock upon order placement.",
        "Media: Support for multiple product images."
    ])

    # 13. Order Processing (Module)
    add_content_slide("Module: Order Processing", [
        "Order Creation: Validates cart items, stock availability, and prices.",
        "Unique Tracking: Generates readable IDs (TANA-YYYYMMDD-XXXX).",
        "Status Workflow: Pending -> Paid -> Approved -> Shipped -> Delivered.",
        "History: Users can view past orders and download summaries."
    ])

    # 14. Payment Integration
    add_content_slide("Payment Integration", [
        "Gateway: Chapa (Ethio-centric payment provider).",
        "Flow: Redirect-based secure checkout.",
        "Webhook Verification: Server-side validation of payment success to prevent fraud.",
        "Currency: Supports local currency (ETB)."
    ])

    # 15. Security Features
    add_content_slide("Security Strategy", [
        "Data Encryption: BCrypt for passwords.",
        "API Security: Rate limiting, CORS configuration, and Input Validation (Express-Validator).",
        "Authorization: Middleware checks for valid Tokens and Roles before access.",
        "Audit Logs: Tracks critical actions (Logins, Deletions) for accountability."
    ])

    # 16. Database Design
    add_content_slide("Database & Persistence", [
        "Collections: Users, Products, Orders, Notifications, Comments.",
        "Relationships: Referenced via MongoDB ObjectIds (e.g., Order belongs to User).",
        "Indexes: Optimized for search performance (Product Name, Category).",
        "Scalability: Schema-less design allows easy future attribute additions."
    ])

    # 17. Challenges & Solutions
    add_content_slide("Challenges & Solutions", [
        "Challenge: Handling concurrent orders for low-stock items.",
        "Solution: Mongoose transactions and atomic updates.",
        "Challenge: Ensuring secure payments without PCI-DSS overhead.",
        "Solution: Offloading sensitive card data handling to Chapa's secure pages."
    ])

    # 18. Conclusion
    add_content_slide("Conclusion & Future Work", [
        "The project successfully met the core objectives of creating a functional e-commerce platform.",
        "Provides a solid foundation for local businesses to digitize.",
        "Future Enhancements:",
        "   - AI-based Product Recommendations.",
        "   - Native Mobile Applications (iOS/Android).",
        "   - Advanced Analytics & Forecasting."
    ])

    # Save
    output_path = "Tana_Market_Premium_Presentation.pptx"
    prs.save(output_path)
    print(f"Successfully generated {output_path}")

if __name__ == "__main__":
    create_presentation()
